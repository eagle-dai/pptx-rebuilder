import os
import io
import re
import base64
import json
import zipfile
from typing import Optional
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from langchain_core.messages import HumanMessage
from PIL import Image

# SAP Cloud SDK for AI - 使用 LangChain 集成
# 文档: https://help.sap.com/doc/generative-ai-hub-sdk/CLOUD/en-US/index.html
from gen_ai_hub.proxy.langchain.init_models import init_llm

app = FastAPI(title="PPTX Rebuilder API")

# 允许前端跨域请求
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # 生产环境请修改为实际前端地址
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# SAP AI Core 凭证说明：
# 需要在环境变量或 ~/.aicore/config.json 中配置以下凭证：
#   - AICORE_CLIENT_ID: SAP AI Core 客户端 ID
#   - AICORE_CLIENT_SECRET: SAP AI Core 客户端密钥
#   - AICORE_AUTH_URL: 认证 URL（例如 https://*.authentication.sap.hana.ondemand.com/oauth/token）
#   - AICORE_BASE_URL: 服务 URL（例如 https://api.ai.*.cfapps.sap.hana.ondemand.com/v2）
#   - AICORE_RESOURCE_GROUP: 资源组名称


def fix_json_quotes(json_str: str) -> str:
    """
    修复 JSON 字符串中的智能引号问题。

    LLM 返回的 JSON 可能：
    1. 使用智能引号（"" ''）作为 JSON 结构引号
    2. 在字符串值内部包含未转义的引号（如 "Fast Brain"）

    修复策略：
    1. 将智能引号转换为标准 ASCII 引号
    2. 智能识别字符串值内的引号并替换为单引号
    """
    # 第一步：将智能单引号统一替换为 ASCII 单引号（这不会破坏 JSON）
    json_str = json_str.replace(''', "'").replace(''', "'")

    # 第二步：将智能双引号替换为标准双引号
    json_str = json_str.replace('"', '"').replace('"', '"')

    # 第三步：处理字符串值内部的双引号
    # 逐字符扫描，识别 JSON 字符串边界，将值内部的双引号替换为单引号
    result = []
    in_string = False
    i = 0

    while i < len(json_str):
        c = json_str[i]

        # 处理转义字符
        if c == '\\' and i + 1 < len(json_str):
            result.append(c)
            result.append(json_str[i + 1])
            i += 2
            continue

        if c == '"':
            if not in_string:
                # 进入字符串
                in_string = True
                result.append(c)
            else:
                # 可能是字符串结束，也可能是值内部的引号
                # 向后看：跳过空白后，如果是 , } ] : 则是字符串结束
                j = i + 1
                while j < len(json_str) and json_str[j] in ' \t\n\r':
                    j += 1

                if j >= len(json_str) or json_str[j] in ',}]:':
                    # 字符串结束
                    in_string = False
                    result.append(c)
                else:
                    # 值内部的引号，替换为单引号
                    result.append("'")
        else:
            result.append(c)

        i += 1

    return ''.join(result)


def get_llm():
    """
    获取 SAP AI SDK 的 Claude 语言模型实例。

    使用 SAP Generative AI Hub 通过 Amazon Bedrock 代理访问 Claude 模型。
    可通过环境变量 LLM_MODEL_NAME 配置模型名称。

    常见的模型名称：
    - anthropic--claude-3-5-sonnet（推荐，支持 Vision）
    - anthropic--claude-3-opus
    - anthropic--claude-3-sonnet
    - anthropic--claude-3-haiku
    """
    model_name = os.getenv("LLM_MODEL_NAME", "anthropic--claude-4.5-sonnet")
    # 只设置 max_tokens，让 SDK 使用默认的采样参数
    llm = init_llm(model_name, max_tokens=8192)
    # 尝试覆盖可能导致冲突的参数
    if hasattr(llm, 'top_p'):
        llm.top_p = None
    return llm


def extract_slide_images_from_pptx(pptx_bytes: io.BytesIO) -> list[tuple[int, str, list[tuple[str, bytes]]]]:
    """
    从 PPTX 文件中提取每张幻灯片的图像及其所有嵌入图片。

    对于"图片型 PPTX"（每张幻灯片主体是一张图片），直接提取嵌入的图片。
    返回: [(slide_index, base64_main_image, [(image_id, image_bytes), ...]), ...]

    - base64_main_image: 幻灯片主图（用于 LLM 分析）
    - embedded_images: 该幻灯片引用的所有图片原始数据（用于重建时插入）

    实现思路：
    1. PPTX 是 ZIP 格式，图片存储在 ppt/media/ 目录下
    2. 幻灯片与图片的对应关系在 ppt/slides/_rels/slideX.xml.rels 中定义
    """
    images = []
    pptx_bytes.seek(0)

    with zipfile.ZipFile(pptx_bytes, 'r') as zf:
        # 获取所有幻灯片文件，按编号排序
        slide_files = sorted([
            f for f in zf.namelist()
            if f.startswith('ppt/slides/slide') and f.endswith('.xml')
        ], key=lambda x: int(x.replace('ppt/slides/slide', '').replace('.xml', '')))

        for slide_idx, slide_file in enumerate(slide_files):
            # 获取该幻灯片对应的关系文件
            slide_num = slide_file.replace('ppt/slides/slide', '').replace('.xml', '')
            rels_file = f'ppt/slides/_rels/slide{slide_num}.xml.rels'

            if rels_file not in zf.namelist():
                continue

            # 解析关系文件，找到图片引用
            rels_content = zf.read(rels_file).decode('utf-8')

            # 查找所有图片引用（格式如 Target="../media/image1.png"）
            image_refs = re.findall(r'Target="\.\./(media/image\d+\.[a-zA-Z]+)"', rels_content)

            # 提取所有图片的原始字节数据（用于重建时插入）
            embedded_images = []
            main_image_base64 = None

            for i, image_ref in enumerate(image_refs):
                image_path = f'ppt/{image_ref}'
                if image_path in zf.namelist():
                    image_data = zf.read(image_path)
                    # 保存原始图片数据
                    image_id = f"image_{slide_idx}_{i}"
                    embedded_images.append((image_id, image_data))

                    # 第一张图片作为主图用于 LLM 分析
                    if i == 0:
                        img = Image.open(io.BytesIO(image_data))
                        png_buffer = io.BytesIO()
                        img.convert('RGB').save(png_buffer, format='PNG')
                        png_buffer.seek(0)
                        main_image_base64 = base64.b64encode(png_buffer.read()).decode('utf-8')

            if main_image_base64:
                images.append((slide_idx, main_image_base64, embedded_images))

    return images


def parse_color(color_str: str) -> Optional[RGBColor]:
    """
    解析颜色字符串，支持多种格式。

    支持格式：
    - "#RRGGBB" 或 "#RGB"（十六进制）
    - "rgb(r, g, b)"
    - 颜色名称（如 "white", "black", "red"）
    """
    if not color_str:
        return None

    color_str = color_str.strip().lower()

    # 颜色名称映射
    color_names = {
        "white": RGBColor(255, 255, 255),
        "black": RGBColor(0, 0, 0),
        "red": RGBColor(255, 0, 0),
        "green": RGBColor(0, 128, 0),
        "blue": RGBColor(0, 0, 255),
        "yellow": RGBColor(255, 255, 0),
        "orange": RGBColor(255, 165, 0),
        "gray": RGBColor(128, 128, 128),
        "grey": RGBColor(128, 128, 128),
        "navy": RGBColor(0, 0, 128),
        "purple": RGBColor(128, 0, 128),
    }

    if color_str in color_names:
        return color_names[color_str]

    # 十六进制格式
    if color_str.startswith("#"):
        hex_color = color_str[1:]
        if len(hex_color) == 3:
            hex_color = ''.join([c * 2 for c in hex_color])
        if len(hex_color) == 6:
            try:
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                return RGBColor(r, g, b)
            except ValueError:
                pass

    # rgb() 格式
    rgb_match = re.match(r'rgb\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)', color_str)
    if rgb_match:
        r, g, b = map(int, rgb_match.groups())
        return RGBColor(min(r, 255), min(g, 255), min(b, 255))

    return None


def create_slide_from_analysis(
    prs: Presentation,
    analysis: dict,
    slide_idx: int,
    embedded_images: list[tuple[str, bytes]] = None
):
    """
    根据 Claude Vision 分析结果创建幻灯片。

    参数:
        prs: python-pptx Presentation 对象
        analysis: Claude 返回的幻灯片分析结果
        slide_idx: 幻灯片索引
        embedded_images: 该幻灯片的嵌入图片列表 [(image_id, image_bytes), ...]

    布局策略:
    - 使用空白布局，手动添加文本框以获得最大灵活性
    - 根据 LLM 返回的位置和样式信息精确放置元素
    - 保留原始图片
    """
    # 使用空白布局（索引 6 通常是空白布局）
    try:
        blank_layout = prs.slide_layouts[6]
    except IndexError:
        blank_layout = prs.slide_layouts[-1]

    slide = prs.slides.add_slide(blank_layout)

    # 幻灯片尺寸
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 获取背景颜色
    bg_color = analysis.get('background_color')
    if bg_color:
        parsed_color = parse_color(bg_color)
        if parsed_color:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = parsed_color

    # 处理所有元素，按 z_order 排序（如果有的话）
    elements = analysis.get('elements', [])

    # 如果没有 elements 字段，使用旧格式兼容
    if not elements:
        elements = _convert_legacy_format_to_elements(analysis, slide_idx)

    # 按 z_order 排序，确保正确的层叠顺序
    elements.sort(key=lambda x: x.get('z_order', 0))

    # 用于追踪图片索引
    image_index = 0

    for element in elements:
        elem_type = element.get('type', 'text')

        # 解析位置（百分比或绝对值）
        left = _parse_position(element.get('left', 0), slide_width)
        top = _parse_position(element.get('top', 0), slide_height)
        width = _parse_position(element.get('width', 100), slide_width)
        height = _parse_position(element.get('height', 50), slide_height)

        if elem_type == 'text':
            _add_text_element(slide, element, left, top, width, height)
        elif elem_type == 'image':
            # 插入实际图片
            if embedded_images and image_index < len(embedded_images):
                _add_image_element(slide, embedded_images[image_index][1], left, top, width, height)
                image_index += 1
        elif elem_type == 'shape':
            _add_shape_element(slide, element, left, top, width, height)
        elif elem_type == 'table':
            _add_table_element(slide, element, left, top, width, height)

    return slide


def _convert_legacy_format_to_elements(analysis: dict, slide_idx: int) -> list:
    """
    将旧格式的分析结果转换为新的 elements 格式，保持向后兼容。
    """
    elements = []
    current_top_pct = 5  # 从顶部 5% 开始

    # 添加标题
    title_text = analysis.get('title', f'幻灯片 {slide_idx + 1}')
    if title_text:
        elements.append({
            'type': 'text',
            'content': title_text,
            'left': '5%',
            'top': f'{current_top_pct}%',
            'width': '90%',
            'height': '12%',
            'font_size': 32,
            'bold': True,
            'z_order': 1
        })
        current_top_pct += 15

    # 添加副标题
    subtitle_text = analysis.get('subtitle')
    if subtitle_text:
        elements.append({
            'type': 'text',
            'content': subtitle_text,
            'left': '5%',
            'top': f'{current_top_pct}%',
            'width': '90%',
            'height': '8%',
            'font_size': 20,
            'italic': True,
            'color': '#646464',
            'z_order': 2
        })
        current_top_pct += 10

    # 添加正文
    body_items = analysis.get('body_items', [])
    if body_items:
        body_text = '\n'.join([
            f"{'    ' * item.get('level', 0)}• {item.get('text', '')}"
            for item in body_items
        ])
        elements.append({
            'type': 'text',
            'content': body_text,
            'left': '5%',
            'top': f'{current_top_pct}%',
            'width': '90%',
            'height': '60%',
            'font_size': 18,
            'z_order': 3
        })

    # 添加图片占位
    images = analysis.get('images', [])
    for i, img_desc in enumerate(images):
        elements.append({
            'type': 'image',
            'description': img_desc if isinstance(img_desc, str) else img_desc.get('description', ''),
            'left': img_desc.get('left', '60%') if isinstance(img_desc, dict) else '60%',
            'top': img_desc.get('top', '30%') if isinstance(img_desc, dict) else '30%',
            'width': img_desc.get('width', '35%') if isinstance(img_desc, dict) else '35%',
            'height': img_desc.get('height', '50%') if isinstance(img_desc, dict) else '50%',
            'z_order': 10 + i
        })

    # 添加页脚
    other_texts = analysis.get('other_texts', [])
    if other_texts:
        elements.append({
            'type': 'text',
            'content': ' | '.join(other_texts),
            'left': '5%',
            'top': '92%',
            'width': '90%',
            'height': '5%',
            'font_size': 10,
            'color': '#808080',
            'align': 'center',
            'z_order': 100
        })

    return elements


def _parse_position(value, reference_size) -> int:
    """
    解析位置值，支持百分比和绝对像素值。

    - "50%" -> reference_size 的 50%
    - 100 -> 100 EMU（python-pptx 内部单位）
    - "100px" -> 约 100 像素（转换为 EMU）
    """
    if isinstance(value, str):
        value = value.strip()
        if value.endswith('%'):
            pct = float(value[:-1]) / 100
            return int(reference_size * pct)
        elif value.endswith('px'):
            # 1 inch ≈ 96 px，1 inch = 914400 EMU
            px = float(value[:-2])
            return int(px * 914400 / 96)
    return int(value) if value else 0


def _add_text_element(slide, element: dict, left: int, top: int, width: int, height: int):
    """添加文本元素到幻灯片。"""
    content = element.get('content', '')
    if not content:
        return

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    # 设置垂直对齐
    v_align = element.get('vertical_align', 'top')
    if v_align == 'middle':
        text_frame.anchor = MSO_ANCHOR.MIDDLE
    elif v_align == 'bottom':
        text_frame.anchor = MSO_ANCHOR.BOTTOM
    else:
        text_frame.anchor = MSO_ANCHOR.TOP

    # 处理多行文本
    lines = content.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            para = text_frame.paragraphs[0]
        else:
            para = text_frame.add_paragraph()

        para.text = line

        # 设置字体属性
        font_size = element.get('font_size', 18)
        para.font.size = Pt(font_size)
        para.font.name = element.get('font_name', 'Microsoft YaHei')

        if element.get('bold'):
            para.font.bold = True
        if element.get('italic'):
            para.font.italic = True

        # 设置颜色
        color = element.get('color')
        if color:
            parsed_color = parse_color(color)
            if parsed_color:
                para.font.color.rgb = parsed_color

        # 设置对齐
        align = element.get('align', 'left')
        if align == 'center':
            para.alignment = PP_ALIGN.CENTER
        elif align == 'right':
            para.alignment = PP_ALIGN.RIGHT
        else:
            para.alignment = PP_ALIGN.LEFT


def _add_image_element(slide, image_data: bytes, left: int, top: int, width: int, height: int):
    """添加图片元素到幻灯片。"""
    try:
        image_stream = io.BytesIO(image_data)
        slide.shapes.add_picture(image_stream, left, top, width, height)
    except Exception as e:
        print(f"添加图片失败: {e}")


def _add_shape_element(slide, element: dict, left: int, top: int, width: int, height: int):
    """添加形状元素到幻灯片（如矩形、圆形等）。"""
    from pptx.enum.shapes import MSO_SHAPE

    shape_type = element.get('shape_type', 'rectangle')
    shape_map = {
        'rectangle': MSO_SHAPE.RECTANGLE,
        'rounded_rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
        'oval': MSO_SHAPE.OVAL,
        'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
    }

    mso_shape = shape_map.get(shape_type, MSO_SHAPE.RECTANGLE)
    shape = slide.shapes.add_shape(mso_shape, left, top, width, height)

    # 设置填充颜色
    fill_color = element.get('fill_color')
    if fill_color:
        parsed_color = parse_color(fill_color)
        if parsed_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = parsed_color

    # 设置边框
    line_color = element.get('line_color')
    if line_color:
        parsed_color = parse_color(line_color)
        if parsed_color:
            shape.line.color.rgb = parsed_color


def _add_table_element(slide, element: dict, left: int, top: int, width: int, height: int):
    """添加表格元素到幻灯片。"""
    rows_data = element.get('rows', [])
    if not rows_data:
        return

    num_rows = len(rows_data)
    num_cols = max(len(row) for row in rows_data) if rows_data else 1

    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    # 填充表格数据
    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            if col_idx < num_cols:
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_text)

                # 设置表头样式
                if row_idx == 0 and element.get('has_header', True):
                    para = cell.text_frame.paragraphs[0]
                    para.font.bold = True
                    para.font.size = Pt(14)


def analyze_slide_image(llm, image_base64: str, slide_index: int) -> dict:
    """
    使用 Claude Vision 分析幻灯片图像，提取布局、文本和样式信息。

    参数:
        llm: SAP AI SDK LangChain LLM 实例
        image_base64: Base64 编码的幻灯片图像
        slide_index: 幻灯片索引（从 0 开始）

    返回:
        包含幻灯片布局、文本和样式信息的字典
    """
    # 构建包含图像的消息（LangChain 格式）
    message = HumanMessage(
        content=[
            {
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": "image/png",
                    "data": image_base64,
                },
            },
            {
                "type": "text",
                "text": """请仔细分析这张幻灯片图片，提取所有可见元素的内容、位置和样式，返回 JSON 格式。

返回格式：
{
  "background_color": "背景色，如 #FFFFFF 或 white（如果是纯色）",
  "elements": [
    {
      "type": "text",
      "content": "文本内容（保留换行符 \\n）",
      "left": "距左边距离（百分比，如 5%）",
      "top": "距顶部距离（百分比，如 10%）",
      "width": "宽度（百分比，如 90%）",
      "height": "高度（百分比，如 15%）",
      "font_size": 32,
      "bold": true,
      "italic": false,
      "color": "文字颜色，如 #333333 或 white",
      "align": "left/center/right",
      "vertical_align": "top/middle/bottom",
      "z_order": 1
    },
    {
      "type": "image",
      "description": "图片内容描述",
      "left": "60%",
      "top": "25%",
      "width": "35%",
      "height": "50%",
      "z_order": 2
    },
    {
      "type": "shape",
      "shape_type": "rectangle/rounded_rectangle/oval/triangle",
      "left": "10%",
      "top": "80%",
      "width": "80%",
      "height": "5%",
      "fill_color": "#0066CC",
      "z_order": 0
    },
    {
      "type": "table",
      "rows": [["表头1", "表头2"], ["数据1", "数据2"]],
      "has_header": true,
      "left": "10%",
      "top": "40%",
      "width": "80%",
      "height": "40%",
      "z_order": 3
    }
  ]
}

重要规则：
1. 位置使用百分比，估算元素在幻灯片中的相对位置
2. z_order 表示层叠顺序，数字越小越在底层（背景形状用 0，文字用更高值）
3. 图片元素只需描述内容，我会自动插入原图
4. 识别所有可见的文字、图片、形状、表格
5. 文字颜色和背景色要准确识别（特别是白色文字在深色背景上）
6. 只返回纯 JSON，不要用 markdown 代码块

示例（简单幻灯片）：
{"background_color": "#1A365D", "elements": [{"type": "text", "content": "项目介绍", "left": "5%", "top": "40%", "width": "90%", "height": "20%", "font_size": 44, "bold": true, "color": "white", "align": "center", "z_order": 1}]}"""
            }
        ]
    )

    # 调用 LLM 进行分析
    response = llm.invoke([message])

    # 解析返回的 JSON
    try:
        response_text = response.content

        # 清理可能存在的 markdown 代码块标记
        response_text = response_text.strip()

        # 使用正则表达式提取 JSON 内容（匹配 ```json...``` 或 ```...```）
        code_block_match = re.search(r'```(?:json)?\s*\n?(.*?)\n?```', response_text, re.DOTALL)
        if code_block_match:
            response_text = code_block_match.group(1).strip()

        # 如果仍然以 ``` 开头，手动处理
        if response_text.startswith("```"):
            lines = response_text.split("\n")
            start_idx = 1
            end_idx = len(lines)
            for i, line in enumerate(lines):
                if i > 0 and line.strip() == "```":
                    end_idx = i
                    break
            response_text = "\n".join(lines[start_idx:end_idx]).strip()

        # 提取 JSON 对象（从第一个 { 到最后一个 }）
        json_start = response_text.find('{')
        json_end = response_text.rfind('}')
        if json_start != -1 and json_end != -1 and json_end > json_start:
            response_text = response_text[json_start:json_end + 1]

        # 修复智能引号问题
        # LLM 返回的 JSON 可能使用智能引号（"" ''）作为 JSON 结构引号
        # 同时字符串值内部也可能包含智能引号（如 "Fast Brain"）
        # 策略：先标准化引号，再处理字符串值内的引号
        response_text = fix_json_quotes(response_text)

        result = json.loads(response_text)

        # 验证返回的结果包含必要字段
        if not isinstance(result, dict):
            raise ValueError("返回结果不是字典类型")

        # 确保有 title 字段
        if 'title' not in result or not result['title']:
            result['title'] = f"幻灯片 {slide_index + 1}"

        return result

    except (json.JSONDecodeError, AttributeError, ValueError) as e:
        # 如果解析失败，记录错误并返回错误提示
        print(f"JSON 解析失败 (幻灯片 {slide_index + 1}): {e}")
        print(f"原始响应: {response.content[:500] if hasattr(response, 'content') else response}...")
        return {
            "title": f"幻灯片 {slide_index + 1}",
            "subtitle": "（内容解析失败，请重试）",
            "body_items": [{"text": "LLM 返回的内容无法解析为有效的 JSON 格式", "level": 0}],
            "images": [],
            "other_texts": []
        }


@app.post("/api/convert")
async def convert_pptx(file: UploadFile = File(...)):
    """
    将图片型 PPTX 转换为可编辑的 PPTX。

    处理流程：
    1. 从上传的 PPTX 中提取每张幻灯片的图片
    2. 使用 Claude Vision 分析每张图片，提取文字和布局
    3. 根据分析结果用 python-pptx 重建幻灯片
    4. 返回新生成的可编辑 PPTX
    """
    if not file.filename.endswith('.pptx'):
        raise HTTPException(status_code=400, detail="仅支持 .pptx 文件上传")

    try:
        # 读取上传的文件
        content = await file.read()
        source_pptx = io.BytesIO(content)

        # 步骤 1: 从 PPTX 中提取幻灯片图片
        slide_images = extract_slide_images_from_pptx(source_pptx)

        if not slide_images:
            raise HTTPException(
                status_code=400,
                detail="未能从 PPTX 中提取到图片，请确保这是一个图片型 PPT 文件"
            )

        # 步骤 2: 初始化 LLM 并分析每张幻灯片
        llm = get_llm()
        slide_analyses = []

        for slide_idx, image_base64, embedded_images in slide_images:
            analysis = analyze_slide_image(llm, image_base64, slide_idx)
            slide_analyses.append((slide_idx, analysis, embedded_images))

        # 步骤 3: 创建新的演示文稿
        target_presentation = Presentation()

        # 设置幻灯片尺寸为 16:9
        target_presentation.slide_width = Inches(13.333)
        target_presentation.slide_height = Inches(7.5)

        # 步骤 4: 根据分析结果创建每张幻灯片（包含原始图片）
        for slide_idx, analysis, embedded_images in slide_analyses:
            create_slide_from_analysis(target_presentation, analysis, slide_idx, embedded_images)

        # 将生成的 PPT 保存到内存中并返回给前端
        output = io.BytesIO()
        target_presentation.save(output)
        output.seek(0)

        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename=rebuilt_{file.filename}"}
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"转换失败: {str(e)}")


@app.get("/api/health")
def health_check():
    """健康检查端点，验证后端服务和 SAP AI SDK 配置状态"""
    # 检查 SAP AI Core 环境变量是否已配置
    required_env_vars = [
        "AICORE_CLIENT_ID",
        "AICORE_CLIENT_SECRET",
        "AICORE_AUTH_URL",
        "AICORE_BASE_URL",
        "AICORE_RESOURCE_GROUP"
    ]

    missing_vars = [var for var in required_env_vars if not os.getenv(var)]

    if missing_vars:
        return {
            "status": "warning",
            "message": "Backend is running but SAP AI Core not fully configured",
            "missing_env_vars": missing_vars
        }

    return {"status": "ok", "message": "Backend is running with SAP AI Core configured"}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
